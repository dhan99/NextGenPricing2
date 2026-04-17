"""Build exports/Integrations-API-Overview.pptx from a structured slide model.

Mirrors docs/integrations/api-overview-slides.md and the narrative in
docs/integrations/api-overview.md. Run:  python scripts/build_integrations_deck.py
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

import sys
sys.path.insert(0, str(Path(__file__).resolve().parent))
from integrations_samples import (
    D365_SAMPLES, D365_PROD_SAMPLES, WD_SAMPLES, WD_PROD_SAMPLES,
)

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "exports" / "Integrations-API-Overview.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

NAVY = RGBColor(0x12, 0x2B, 0x4A)
ACCENT = RGBColor(0xC2, 0x8A, 0x2A)
INK = RGBColor(0x1F, 0x24, 0x2E)
MUTED = RGBColor(0x5C, 0x66, 0x73)
LIGHT = RGBColor(0xF5, 0xF1, 0xE8)
LINE = RGBColor(0xD9, 0xD2, 0xC2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


def add_blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def text_box(slide, x, y, w, h, text, *, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        run.font.name = font; run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color
    return tb


def header(slide, eyebrow, title, *, footer_idx=None, footer_total=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
    fill(bar, ACCENT)
    text_box(slide, Inches(0.55), Inches(0.30), Inches(10), Inches(0.3),
             eyebrow.upper(), size=10, bold=True, color=ACCENT)
    text_box(slide, Inches(0.55), Inches(0.55), Inches(12), Inches(0.7),
             title, size=26, bold=True, color=NAVY)
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0.55), Inches(1.25), Inches(1.4), Emu(20000))
    fill(underline, NAVY)
    if footer_idx and footer_total:
        text_box(slide, Inches(11.7), Inches(7.05), Inches(1.5), Inches(0.3),
                 f"{footer_idx} / {footer_total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    text_box(slide, Inches(0.55), Inches(7.05), Inches(8), Inches(0.3),
             "DealPad · Integrations API Overview · April 2026", size=9, color=MUTED)


# ---------- slide builders ----------

def slide_title(prs, total):
    s = add_blank(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H); fill(bg, NAVY)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.3), SLIDE_W, Inches(0.06))
    fill(accent, ACCENT)
    text_box(s, Inches(0.8), Inches(2.4), Inches(11), Inches(0.5),
             "DEALPAD · STAKEHOLDER BRIEFING", size=14, bold=True, color=ACCENT)
    text_box(s, Inches(0.8), Inches(2.9), Inches(12), Inches(1.2),
             "Integrations API Overview", size=44, bold=True, color=WHITE)
    text_box(s, Inches(0.8), Inches(4.6), Inches(12), Inches(0.8),
             "Microsoft Dynamics 365 (CRM)  ·  Workday (HCM / Financial Management)",
             size=20, color=WHITE)
    text_box(s, Inches(0.8), Inches(6.7), Inches(12), Inches(0.4),
             "April 2026  ·  Persistent simulation today, configuration-only cutover to live",
             size=12, color=ACCENT)


def slide_section(prs, label, idx, total):
    s = add_blank(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H); fill(bg, LIGHT)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8), SLIDE_W, Inches(1.9))
    fill(band, NAVY)
    text_box(s, Inches(0.8), Inches(3.0), Inches(12), Inches(0.5),
             "SECTION", size=12, bold=True, color=ACCENT)
    text_box(s, Inches(0.8), Inches(3.4), Inches(12), Inches(1.2),
             label, size=40, bold=True, color=WHITE)
    text_box(s, Inches(11.7), Inches(7.05), Inches(1.5), Inches(0.3),
             f"{idx} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def bullets(slide, x, y, w, h, items, *, size=14, color=INK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run(); run.text = f"•  {it}"
        run.font.name = "Calibri"; run.font.size = Pt(size); run.font.color.rgb = color


def slide_bullets(prs, eyebrow, title, items, idx, total, *, intro=None):
    s = add_blank(prs); header(s, eyebrow, title, footer_idx=idx, footer_total=total)
    y = Inches(1.7)
    if intro:
        text_box(s, Inches(0.55), y, Inches(12), Inches(0.6), intro, size=14, color=MUTED)
        y = Inches(2.3)
    bullets(s, Inches(0.55), y, Inches(12.3), Inches(5), items, size=16)


def table(slide, x, y, w, h, headers, rows, *, header_fill=NAVY, header_color=WHITE,
          col_widths=None, font_size=10, header_size=10):
    tbl_shape = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h)
    tbl = tbl_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    # header row
    for i, htext in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = htext
        run.font.size = Pt(header_size); run.font.bold = True
        run.font.color.rgb = header_color; run.font.name = "Calibri"
        cell.margin_left = cell.margin_right = Inches(0.06)
        cell.margin_top = cell.margin_bottom = Inches(0.03)
    # data rows
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run(); run.text = str(val)
            run.font.size = Pt(font_size); run.font.name = "Calibri"
            run.font.color.rgb = INK
            if c == 0:
                run.font.name = "Consolas"; run.font.size = Pt(font_size - 1)
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
    return tbl


def code_block(slide, x, y, w, h, text, *, label=None):
    if label:
        text_box(slide, x, y - Inches(0.3), w, Inches(0.3),
                 label, size=11, bold=True, color=NAVY)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1E, 0x26, 0x33)
    box.line.color.rgb = NAVY
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = Inches(0.08)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run(); run.text = line if line else " "
        run.font.name = "Consolas"; run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0xEC, 0xE6, 0xD0)


def slide_dataflow(prs, eyebrow, title, idx, total, *, system_name, inbound, outbound, notes):
    s = add_blank(prs); header(s, eyebrow, title, footer_idx=idx, footer_total=total)
    # Two boxes + arrows
    left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.7), Inches(2.1), Inches(4.3), Inches(2.6))
    fill(left, NAVY)
    text_box(s, Inches(0.85), Inches(2.3), Inches(4), Inches(0.5),
             "DealPad", size=20, bold=True, color=WHITE)
    text_box(s, Inches(0.85), Inches(2.85), Inches(4), Inches(2),
             "Node + Postgres\nProvider pattern\nAudit log + sync log\nRBAC + override flow",
             size=12, color=ACCENT)
    right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(8.4), Inches(2.1), Inches(4.3), Inches(2.6))
    fill(right, ACCENT)
    text_box(s, Inches(8.55), Inches(2.3), Inches(4), Inches(0.5),
             system_name, size=20, bold=True, color=NAVY)
    text_box(s, Inches(8.55), Inches(2.85), Inches(4), Inches(2),
             "\n".join(notes), size=12, color=NAVY)
    # Inbound arrow (right -> left)
    a1 = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW,
                            Inches(5.15), Inches(2.5), Inches(3.1), Inches(0.55))
    fill(a1, MUTED)
    text_box(s, Inches(5.15), Inches(2.18), Inches(3.1), Inches(0.3),
             inbound, size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    # Outbound arrow (left -> right)
    a2 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                            Inches(5.15), Inches(3.7), Inches(3.1), Inches(0.55))
    fill(a2, NAVY)
    text_box(s, Inches(5.15), Inches(4.3), Inches(3.1), Inches(0.3),
             outbound, size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    # Bottom callouts
    text_box(s, Inches(0.55), Inches(5.2), Inches(12.3), Inches(0.4),
             "Trigger points", size=14, bold=True, color=NAVY)
    bullets(s, Inches(0.55), Inches(5.55), Inches(12.3), Inches(1.5),
            triggers_for(system_name), size=13)


def triggers_for(name):
    if "Dynamics" in name:
        return [
            "Outbound auto-push on deal status / stage / step / fee / cost / hours / margin change (per-trigger toggles in dynamics_settings).",
            "Inbound nightly batch (POST /api/dynamics/nightly-batch) and on-demand pull (POST /api/dynamics/sync).",
            "Manual push: POST /api/dynamics/deals/:id/push.",
        ]
    return [
        "Auto-validate on deal save (workday_settings.autoValidateOnSave).",
        "Submission gate: PATCH /api/deals/:id with status=submitted triggers validation; over_budget or staffing_shortfall blocks unless overridden.",
        "Override: POST /api/workday/validations/:id/override — Finance / Service Line Lead only, justification required.",
    ]


# ---------- endpoint inventories ----------

D365_READ = [
    ("GET",   "/api/dynamics/accounts",                     "List all client accounts"),
    ("GET",   "/api/dynamics/accounts/:id",                 "Single account detail"),
    ("GET",   "/api/dynamics/opportunities",                "List opportunities"),
    ("GET",   "/api/dynamics/opportunities/eligible",       "Develop/Propose opps not yet linked"),
    ("GET",   "/api/dynamics/scope-templates",              "Service-line scope templates"),
    ("GET",   "/api/dynamics/pipeline",                     "Pipeline rollup (stage / owner / forecast)"),
    ("GET",   "/api/dynamics/sync-log",                     "Last 100 sync events"),
    ("GET",   "/api/dynamics/settings",                     "Sync toggles"),
    ("GET",   "/api/dynamics/owners",                       "Sales owners + quotas"),
]

D365_WRITE = [
    ("POST",  "/api/dynamics/opportunities",                "Create new opportunity"),
    ("PATCH", "/api/dynamics/opportunities/:id",            "Edit stage / value / owner"),
    ("POST",  "/api/dynamics/opportunities/:id/import",     "Pull opp into DealPad as draft deal"),
    ("POST",  "/api/dynamics/opportunities/:id/unlink",     "Unlink opp from DealPad deal"),
    ("POST",  "/api/dynamics/deals/:id/push",               "Manual push: deal → D365"),
    ("POST",  "/api/dynamics/sync",                         "Bulk on-demand pull / push"),
    ("POST",  "/api/dynamics/nightly-batch",                "Scheduled full sync"),
    ("PATCH", "/api/dynamics/settings",                     "Update sync toggles"),
    ("PATCH", "/api/dynamics/accounts/:id",                 "Edit account record"),
]

D365_MAP = [
    ("GET /api/dynamics/accounts",
     "GET /accounts?$select=name,industrycode,revenue,numberofemployees"),
    ("GET /api/dynamics/accounts/:id",
     "GET /accounts({accountid})"),
    ("GET /api/dynamics/opportunities",
     "GET /opportunities?$expand=parentaccountid($select=name)"),
    ("GET /api/dynamics/opportunities/eligible",
     "GET /opportunities?$filter=statecode eq 0 and stepname in ('Develop','Propose')"),
    ("GET /api/dynamics/pipeline",
     "GET /opportunities?$filter=statecode eq 0&$select=estimatedvalue,closeprobability,stepname"),
    ("GET /api/dynamics/owners",
     "GET /systemusers?$filter=isdisabled eq false"),
    ("POST /api/dynamics/opportunities",
     "POST /opportunities"),
    ("PATCH /api/dynamics/opportunities/:id",
     "PATCH /opportunities({opportunityid})"),
    ("POST /api/dynamics/opportunities/:id/import",
     "GET /opportunities({opportunityid})  +  DealPad insert"),
    ("POST /api/dynamics/deals/:id/push",
     "PATCH /opportunities({opportunityid})"),
    ("POST /api/dynamics/sync",
     "Multiple GET / PATCH on /accounts and /opportunities"),
    ("POST /api/dynamics/nightly-batch",
     "Multiple GET / PATCH on /accounts and /opportunities (use $batch)"),
    ("PATCH /api/dynamics/accounts/:id",
     "PATCH /accounts({accountid})"),
]

WD_READ = [
    ("GET",   "/api/workday/settings",                      "Mode, tenant, tolerances"),
    ("GET",   "/api/workday/cost-centers",                  "Budgets + headroom"),
    ("GET",   "/api/workday/workers",                       "Worker pool + availability"),
    ("GET",   "/api/workday/rate-card",                     "Standard cost rates by role"),
    ("GET",   "/api/workday/validations",                   "Recent validation runs"),
    ("GET",   "/api/workday/validations/:id",               "Validation detail + findings"),
    ("GET",   "/api/workday/deals/:dealId/latest",          "Latest validation for a deal"),
    ("GET",   "/api/workday/events",                        "Last 150 audit events"),
    ("GET",   "/api/workday/dashboard",                     "Cross-deal validation rollup"),
]

WD_WRITE = [
    ("PATCH",  "/api/workday/settings",                     "Update mode / tolerances / credentials"),
    ("POST",   "/api/workday/cost-centers",                 "Create cost center"),
    ("PATCH",  "/api/workday/cost-centers/:id",             "Edit cost center"),
    ("DELETE", "/api/workday/cost-centers/:id",             "Delete cost center"),
    ("POST",   "/api/workday/workers",                      "Create worker"),
    ("PATCH",  "/api/workday/workers/:id",                  "Edit worker"),
    ("DELETE", "/api/workday/workers/:id",                  "Remove worker"),
    ("PATCH",  "/api/workday/rate-card/:id",                "Update standard cost rate"),
    ("POST",   "/api/workday/deals/:dealId/validate",       "Run validation for a deal"),
    ("POST",   "/api/workday/deals/:dealId/link",           "Link / unlink deal ↔ cost center"),
    ("POST",   "/api/workday/validations/:id/override",     "Override blocking validation"),
]

WD_MAP = [
    ("GET /api/workday/cost-centers",
     "GET /financialManagement/v1/{tenant}/costCenters"),
    ("GET /api/workday/workers",
     "GET /staffing/v6/{tenant}/workers"),
    ("GET /api/workday/rate-card",
     "GET /compensation/v1/{tenant}/compensationPlans  (or RaaS report)"),
    ("POST /api/workday/cost-centers",
     "POST /financialManagement/v1/{tenant}/costCenters"),
    ("PATCH /api/workday/cost-centers/:id",
     "PATCH /financialManagement/v1/{tenant}/costCenters/{id}"),
    ("DELETE /api/workday/cost-centers/:id",
     "DELETE /financialManagement/v1/{tenant}/costCenters/{id}"),
    ("POST /api/workday/workers",
     "SOAP Hire_Employee  (Staffing v40+)"),
    ("PATCH /api/workday/workers/:id",
     "SOAP Edit_Position / Change_Job"),
    ("DELETE /api/workday/workers/:id",
     "SOAP Terminate_Employee"),
    ("PATCH /api/workday/rate-card/:id",
     "SOAP Put_Compensation_Plan"),
    ("POST /api/workday/deals/:dealId/validate",
     "Composite: GET costCenters + GET workers + DealPad rules engine"),
    ("POST /api/workday/deals/:dealId/link",
     "DealPad-internal mapping"),
    ("POST /api/workday/validations/:id/override",
     "DealPad-internal audit (workday_validations override fields)"),
]


# ---------- top-level deck assembly ----------

def sample_card(slide, x, y, w, h, method, path, purpose, request, response):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    bg.line.color.rgb = LINE
    badge_w = Inches(0.7)
    badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   x + Inches(0.1), y + Inches(0.1), badge_w, Inches(0.28))
    fill(badge, NAVY)
    tf = badge.text_frame; tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = method
    r.font.name = "Calibri"; r.font.size = Pt(8); r.font.bold = True; r.font.color.rgb = WHITE
    text_box(slide, x + Inches(0.85), y + Inches(0.07), w - Inches(0.95), Inches(0.32),
             path, size=10, bold=True, color=NAVY, font="Consolas")
    text_box(slide, x + Inches(0.1), y + Inches(0.42), w - Inches(0.2), Inches(0.28),
             purpose, size=9, color=MUTED)
    body_y = y + Inches(0.74)
    body_h = (h - Inches(0.84)) / 2
    half_w = w - Inches(0.2)
    code_block(slide, x + Inches(0.1), body_y, half_w, body_h - Inches(0.05),
               request)
    code_block(slide, x + Inches(0.1), body_y + body_h + Inches(0.05),
               half_w, body_h - Inches(0.05), response)


def slide_sample_cards(prs, eyebrow, title, samples, idx, total):
    s = add_blank(prs); header(s, eyebrow, title, footer_idx=idx, footer_total=total)
    # 2x2 grid in usable area below header
    margin = Inches(0.4)
    top = Inches(1.55)
    grid_w = SLIDE_W - 2 * margin
    grid_h = Inches(5.4)
    card_w = (grid_w - Inches(0.3)) / 2
    card_h = (grid_h - Inches(0.3)) / 2
    coords = [
        (margin,                      top),
        (margin + card_w + Inches(0.3), top),
        (margin,                      top + card_h + Inches(0.3)),
        (margin + card_w + Inches(0.3), top + card_h + Inches(0.3)),
    ]
    for sample, (cx, cy) in zip(samples, coords):
        method, path, purpose, req, resp = sample
        sample_card(s, cx, cy, card_w, card_h, method, path, purpose, req, resp)


def slide_prod_samples(prs, eyebrow, title, samples, start_idx, total, *, intro=None):
    """Each entry: (internal_route_label, prod_request_text, prod_response_text).
    Paginates 4 cards per slide; returns the number of slides emitted."""
    pages = [samples[i:i+4] for i in range(0, len(samples), 4)] or [[]]
    for pi, page in enumerate(pages):
        s = add_blank(prs)
        page_title = title if len(pages) == 1 else f"{title} ({pi+1}/{len(pages)})"
        header(s, eyebrow, page_title, footer_idx=start_idx + pi, footer_total=total)
        if intro:
            text_box(s, Inches(0.55), Inches(1.45), Inches(12.3), Inches(0.3),
                     intro, size=10, color=MUTED)
        margin = Inches(0.4)
        top = Inches(1.85)
        grid_w = SLIDE_W - 2 * margin
        grid_h = Inches(5.1)
        card_w = (grid_w - Inches(0.3)) / 2
        card_h = (grid_h - Inches(0.3)) / 2
        coords = [
            (margin,                        top),
            (margin + card_w + Inches(0.3), top),
            (margin,                        top + card_h + Inches(0.3)),
            (margin + card_w + Inches(0.3), top + card_h + Inches(0.3)),
        ]
        for (label, req, resp), (cx, cy) in zip(page, coords):
            bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, card_w, card_h)
            bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.color.rgb = LINE
            text_box(s, cx + Inches(0.1), cy + Inches(0.08), card_w - Inches(0.2), Inches(0.32),
                     label, size=10, bold=True, color=NAVY, font="Consolas")
            body_y = cy + Inches(0.5)
            body_h = (card_h - Inches(0.6)) / 2
            half_w = card_w - Inches(0.2)
            code_block(s, cx + Inches(0.1), body_y, half_w, body_h - Inches(0.05), req)
            code_block(s, cx + Inches(0.1), body_y + body_h + Inches(0.05),
                       half_w, body_h - Inches(0.05), resp)
    return len(pages)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    total = 34

    # 1
    slide_title(prs, total)

    # 2 Executive summary
    s = add_blank(prs); header(s, "Executive Summary", "Why these two integrations matter",
                                footer_idx=2, footer_total=total)
    bullets(s, Inches(0.55), Inches(1.7), Inches(12.3), Inches(5), [
        "Dynamics 365 is the system of record for client accounts and the opportunity pipeline; DealPad pushes fee, stage, probability, and forecast on every deal change.",
        "Workday is the source of truth for budgets, worker availability, and standard cost rates; DealPad gates deal submission on Workday's verdict.",
        "Both integrations run today as persistent simulations behind a Provider interface — full UX, audit, and override flows are real.",
        "Cutover is configuration-only: implement Live providers behind the existing routes, populate secrets, flip mode flag per environment.",
        "Audit is built in: dynamics_sync_log, workday_events, workday_validations + findings capture every call, every override, every actor.",
    ], size=16)

    # 3 D365 section divider
    slide_section(prs, "Microsoft Dynamics 365", 3, total)

    # 4 D365 architecture
    slide_dataflow(prs, "Dynamics 365 · Architecture", "Bi-directional CRM sync",
                   4, total,
                   system_name="Microsoft Dynamics 365",
                   inbound="Inbound: nightly batch + on-demand pull",
                   outbound="Outbound: stage / fee / probability / forecast",
                   notes=["Dataverse Web API v9.2",
                          "Entities: accounts, opportunities,",
                          "  systemusers, contacts",
                          "Native $batch + OData filtering"])

    # 5 D365 auth
    s = add_blank(prs); header(s, "Dynamics 365 · Auth & Security",
                                "OAuth 2.0 client-credentials via Azure AD",
                                footer_idx=5, footer_total=total)
    bullets(s, Inches(0.55), Inches(1.7), Inches(12.3), Inches(2.5), [
        "Token endpoint: https://login.microsoftonline.com/{tenantId}/oauth2/v2.0/token",
        "Scope: https://{org}.api.crm.dynamics.com/.default",
        "Header on every Web API call: Authorization: Bearer <token>",
        "Secrets in Replit Secrets: D365_TENANT_ID, D365_CLIENT_ID, D365_CLIENT_SECRET, D365_ORG_URL",
        "Audit: every inbound / outbound call appended to dynamics_sync_log (direction, entity, fields, actor, trigger, status).",
    ], size=14)
    code_block(s, Inches(0.55), Inches(5.0), Inches(12.3), Inches(2.0),
               "POST /{tenantId}/oauth2/v2.0/token  HTTP/1.1\n"
               "Host: login.microsoftonline.com\n"
               "Content-Type: application/x-www-form-urlencoded\n\n"
               "grant_type=client_credentials\n"
               "&client_id=$D365_CLIENT_ID\n"
               "&client_secret=$D365_CLIENT_SECRET\n"
               "&scope=https://armanino.api.crm.dynamics.com/.default",
               label="Token request")

    # 6 D365 read endpoints
    s = add_blank(prs); header(s, "Dynamics 365 · Endpoints", "Internal read endpoints",
                                footer_idx=6, footer_total=total)
    table(s, Inches(0.55), Inches(1.65), Inches(12.3), Inches(5.2),
          ["Method", "Path", "Purpose"], D365_READ,
          col_widths=[Inches(1.0), Inches(5.0), Inches(6.3)])

    # 7 D365 write endpoints
    s = add_blank(prs); header(s, "Dynamics 365 · Endpoints", "Internal write endpoints",
                                footer_idx=7, footer_total=total)
    table(s, Inches(0.55), Inches(1.65), Inches(12.3), Inches(5.2),
          ["Method", "Path", "Purpose"], D365_WRITE,
          col_widths=[Inches(1.0), Inches(5.0), Inches(6.3)])

    # 8 D365 mapping
    s = add_blank(prs); header(s, "Dynamics 365 · Production Mapping",
                                "Internal route → Dataverse Web API v9.2",
                                footer_idx=8, footer_total=total)
    text_box(s, Inches(0.55), Inches(1.45), Inches(12.3), Inches(0.3),
             "Base: https://{org}.api.crm.dynamics.com/api/data/v9.2/",
             size=11, color=MUTED)
    table(s, Inches(0.55), Inches(1.85), Inches(12.3), Inches(5.0),
          ["DealPad endpoint", "Production Dataverse call"], D365_MAP,
          col_widths=[Inches(5.0), Inches(7.3)], font_size=10)

    # 9 D365 sample
    s = add_blank(prs); header(s, "Dynamics 365 · Sample Payload",
                                "Create opportunity — DealPad ↔ Dataverse",
                                footer_idx=9, footer_total=total)
    code_block(s, Inches(0.55), Inches(2.0), Inches(6.0), Inches(2.4),
               "POST /api/dynamics/opportunities\n"
               "Content-Type: application/json\n\n"
               "{\n"
               "  \"accountId\": 42,\n"
               "  \"name\": \"Crestwood - Annual Audit\",\n"
               "  \"estimatedValue\": 412000,\n"
               "  \"stage\": \"Qualify\",\n"
               "  \"estimatedCloseDate\": \"2026-11-01\",\n"
               "  \"ownerName\": \"Priya Anand\"\n"
               "}",
               label="DealPad request")
    code_block(s, Inches(6.85), Inches(2.0), Inches(6.0), Inches(2.4),
               "POST /api/data/v9.2/opportunities\n"
               "Authorization: Bearer eyJ0eXA...\n"
               "OData-Version: 4.0\n\n"
               "{\n"
               "  \"name\": \"Crestwood - Annual Audit\",\n"
               "  \"estimatedvalue\": 412000,\n"
               "  \"estimatedclosedate\": \"2026-11-01\",\n"
               "  \"stepname\": \"Qualify\",\n"
               "  \"parentaccountid@odata.bind\": \"/accounts(8b3a)\",\n"
               "  \"ownerid@odata.bind\": \"/systemusers(5d7c)\"\n"
               "}",
               label="Dataverse Web API equivalent")
    code_block(s, Inches(0.55), Inches(4.95), Inches(12.3), Inches(2.0),
               "{\n"
               "  \"id\": 137, \"opportunityNumber\": \"OPP-100204\",\n"
               "  \"name\": \"Crestwood - Annual Audit\",\n"
               "  \"estimatedValue\": 412000, \"stage\": \"Qualify\", \"probability\": 20,\n"
               "  \"forecastCategory\": \"Pipeline\", \"syncStatus\": \"queued\",\n"
               "  \"syncDirection\": \"inbound\"\n"
               "}",
               label="Response (DealPad-shaped)")

    # 10 Workday section divider
    slide_section(prs, "Workday", 10, total)

    # 11 Workday architecture
    slide_dataflow(prs, "Workday · Architecture", "Validation gate on save & submit",
                   11, total,
                   system_name="Workday",
                   inbound="Inbound: cost centers, workers, rate card",
                   outbound="Outbound: validation request (composite)",
                   notes=["REST + SOAP web services",
                          "Modules: Financial Management,",
                          "  Staffing, Compensation",
                          "Validation runs gate submission"])

    # 12 Workday auth
    s = add_blank(prs); header(s, "Workday · Auth & Security",
                                "OAuth 2.0 (REST) or ISU + Basic auth (SOAP)",
                                footer_idx=12, footer_total=total)
    bullets(s, Inches(0.55), Inches(1.7), Inches(12.3), Inches(2.5), [
        "OAuth 2.0 token endpoint: https://{host}.workday.com/ccx/oauth2/{tenant}/token",
        "Legacy SOAP services: Authorization: Basic base64(ISU@tenant:password)",
        "Secrets in workday_settings (encrypted): tenantUrl, isuUsername, apiClientId, apiClientSecret",
        "TLS 1.2+ + IP allow-list on Workday tenant; integration system limited via domain security policy",
        "Audit: workday_events (every call), workday_validations + workday_validation_findings (every check + override)",
    ], size=14)
    code_block(s, Inches(0.55), Inches(5.0), Inches(12.3), Inches(2.0),
               "POST /ccx/oauth2/armanino/token  HTTP/1.1\n"
               "Host: wd5.workday.com\n"
               "Content-Type: application/x-www-form-urlencoded\n"
               "Authorization: Basic <base64(client_id:client_secret)>\n\n"
               "grant_type=client_credentials\n"
               "&scope=staffing financialManagement compensation",
               label="OAuth token request")

    # 13 WD read endpoints
    s = add_blank(prs); header(s, "Workday · Endpoints", "Internal read endpoints",
                                footer_idx=13, footer_total=total)
    table(s, Inches(0.55), Inches(1.65), Inches(12.3), Inches(5.2),
          ["Method", "Path", "Purpose"], WD_READ,
          col_widths=[Inches(1.0), Inches(5.0), Inches(6.3)])

    # 14 WD write endpoints
    s = add_blank(prs); header(s, "Workday · Endpoints", "Internal write endpoints",
                                footer_idx=14, footer_total=total)
    table(s, Inches(0.55), Inches(1.65), Inches(12.3), Inches(5.2),
          ["Method", "Path", "Purpose"], WD_WRITE,
          col_widths=[Inches(1.0), Inches(5.0), Inches(6.3)])

    # 15 WD mapping
    s = add_blank(prs); header(s, "Workday · Production Mapping",
                                "Internal route → Workday REST or SOAP",
                                footer_idx=15, footer_total=total)
    text_box(s, Inches(0.55), Inches(1.45), Inches(12.3), Inches(0.3),
             "REST base: https://{host}.workday.com/ccx/api/{service}/v{n}/{tenant}/  ·  "
             "SOAP base: https://{host}.workday.com/ccx/service/{tenant}/{service}/v{n}",
             size=10, color=MUTED)
    table(s, Inches(0.55), Inches(1.85), Inches(12.3), Inches(5.0),
          ["DealPad endpoint", "Production Workday call"], WD_MAP,
          col_widths=[Inches(5.0), Inches(7.3)], font_size=10)

    # 16 WD samples (REST + SOAP)
    s = add_blank(prs); header(s, "Workday · Sample Payload",
                                "Validate deal — REST cost-center pull + SOAP rate-card update",
                                footer_idx=16, footer_total=total)
    code_block(s, Inches(0.55), Inches(2.0), Inches(6.0), Inches(2.4),
               "GET /ccx/api/financialManagement/v1/\n"
               "    armanino/costCenters/CC-CONS-300\n"
               "Authorization: Bearer eyJraWQiOi...\n"
               "Accept: application/json\n\n"
               "200 OK\n"
               "{ \"id\": \"8e1b\", \"code\": \"CC-CONS-300\",\n"
               "  \"name\": \"Technology Consulting\",\n"
               "  \"totalBudget\": 6200000,\n"
               "  \"committed\": 5950000 }",
               label="REST: cost-center read")
    code_block(s, Inches(6.85), Inches(2.0), Inches(6.0), Inches(2.4),
               "<env:Envelope ...>\n"
               " <env:Body>\n"
               "  <wd:Put_Compensation_Plan_Request>\n"
               "   <wd:Plan_Reference Descriptor=\"Senior Manager\"/>\n"
               "   <wd:Plan_Data>\n"
               "    <wd:Standard_Hourly_Cost_Rate>200</wd:...>\n"
               "    <wd:Effective_Date>2026-04-17</wd:...>\n"
               "   </wd:Plan_Data>\n"
               "  </wd:Put_Compensation_Plan_Request>\n"
               " </env:Body>\n</env:Envelope>",
               label="SOAP: rate-card update")
    code_block(s, Inches(0.55), Inches(4.95), Inches(12.3), Inches(2.0),
               "{\n"
               "  \"ok\": false, \"status\": \"staffing_shortfall\", \"validationId\": 412,\n"
               "  \"summary\": \"Staffing shortfall: 240h across roles.\",\n"
               "  \"findings\": [\n"
               "    { \"findingType\": \"budget\", \"severity\": \"info\", \"message\": \"...\" },\n"
               "    { \"findingType\": \"staffing\", \"severity\": \"blocker\", \"roleName\": \"Sr Consultant\",\n"
               "      \"requiredHours\": 640, \"availableHours\": 400, \"shortfallHours\": 240 }\n"
               "  ]\n"
               "}",
               label="Response (DealPad validation)")

    # 17 D365 field map
    s = add_blank(prs); header(s, "Field Mapping · Dynamics 365",
                                "DealPad Deal ↔ D365 Opportunity",
                                footer_idx=17, footer_total=total)
    table(s, Inches(0.55), Inches(1.65), Inches(12.3), Inches(5.2),
          ["DealPad", "Dynamics 365", "Notes"], [
            ("deals.title",                 "name",                        "1:1"),
            ("deals.totalFee",              "estimatedvalue / actualvalue","Won → also actualclosedate"),
            ("deals.endDate",               "estimatedclosedate",          ""),
            ("derived stage",               "stepname",                    "won/lost/approved/submitted/in_review → Won/Lost/Close/Propose"),
            ("derived probability",         "closeprobability",            "Qualify 20 / Develop 40 / Propose 65 / Close 85 / Won 100 / Lost 0"),
            ("derived forecastCategory",    "forecastcategory",            "Won/Lost→Closed, ≥80→Commit, ≥50→Best Case, else Pipeline"),
            ("deals.pdlName",               "ownerid (systemusers)",       "Resolved via dynamics_owners"),
            ("clients.name",                "parentaccountid (accounts)",  ""),
          ],
          col_widths=[Inches(3.5), Inches(4.5), Inches(4.3)], font_size=10)

    # 18 Workday field map
    s = add_blank(prs); header(s, "Field Mapping · Workday",
                                "DealPad pricing line ↔ Workday worker availability + cost center",
                                footer_idx=18, footer_total=total)
    table(s, Inches(0.55), Inches(1.65), Inches(12.3), Inches(5.2),
          ["DealPad", "Workday", "Notes"], [
            ("pricing_lines.roleId → roles.name", "compensationPlans.role",                "Drives rate-card lookup"),
            ("pricing_lines.costRate",            "compensationPlans.standardHourlyCostRate","Variance > tolerance → rate_variance"),
            ("Σ pricing_lines.hours by role",     "Σ workers.availableHours by role",       "Required vs available → staffing_shortfall"),
            ("deals.totalCost",                   "costCenters.committed (delta)",          "Pre-commit headroom → over_budget"),
            ("deals.workdayCostCenterId",         "costCenters.id",                          "One cost center per deal"),
            ("deals.businessUnit",                "costCenters.businessUnit",               "Default mapping when no explicit link"),
          ],
          col_widths=[Inches(3.8), Inches(4.5), Inches(4.0)], font_size=10)

    # 19 Risks & rate limits
    s = add_blank(prs); header(s, "Operations", "Risks, rate limits, and audit",
                                footer_idx=19, footer_total=total)
    bullets(s, Inches(0.55), Inches(1.7), Inches(12.3), Inches(5), [
        "Dynamics 365 Web API: 6,000 requests / 5-minute sliding window per user — use $batch for nightly job; honor Retry-After on 429.",
        "Workday REST: ~1,500 req/min per tenant; SOAP services serialize per ISU — keep concurrency ≤ 8 for Get_Workers.",
        "Network: Workday IP allow-list must include production egress; Dynamics requires Azure AD tenant trust granted by IT.",
        "Data drift: simulated tables flushed (or migrated) at cutover so we don't serve stale rows after switching mode → live.",
        "Override audit: blocking Workday validations only overridable by Finance / Service Line Lead, justification ≥ 5 chars, captured in workday_validations.overrideJustification + overriddenBy.",
        "Audit log retention: dynamics_sync_log + workday_events queryable from the UI; export pipeline TBD with Security.",
    ], size=14)

    # 20 Cutover checklist
    s = add_blank(prs); header(s, "Cutover", "Sandbox → live in 8 steps",
                                footer_idx=20, footer_total=total)
    items = [
        "Provision sandbox tenants — Dynamics 365 (Dataverse) + Workday Implementation tenant.",
        "Create Azure AD app registration; grant Dataverse user_impersonation scope.",
        "Create Workday ISU + integration system; grant Get/Put domain security on Staffing, Financial Management, Compensation.",
        "Populate Replit Secrets: D365_* and WORKDAY_* keys.",
        "Implement Live providers (LiveDynamicsProvider, LiveWorkdayProvider) behind the existing Provider interface — no route changes.",
        "Flip dynamics_settings.mode / workday_settings.mode to live per environment.",
        "Run shadow-mode for 1 week (read live, write to sim) and reconcile via dynamics_sync_log + workday_events.",
        "Cut over writes; archive simulated tables.",
    ]
    tb = s.shapes.add_textbox(Inches(0.55), Inches(1.7), Inches(12.3), Inches(5.2))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run(); run.text = f"{i+1}.  {it}"
        run.font.name = "Calibri"; run.font.size = Pt(15); run.font.color.rgb = INK

    # ---------- Appendix: per-endpoint sample request + response ----------
    cursor = 21
    slide_section(prs, "Appendix · Per-Endpoint Samples", cursor, total); cursor += 1

    d365_chunks = [D365_SAMPLES[i:i+4] for i in range(0, len(D365_SAMPLES), 4)]
    for i, chunk in enumerate(d365_chunks, start=1):
        slide_sample_cards(
            prs, "Dynamics 365 · Samples",
            f"Internal endpoint request + response ({i}/{len(d365_chunks)})",
            chunk, cursor, total)
        cursor += 1

    cursor += slide_prod_samples(
        prs, "Dynamics 365 · Production Samples",
        "Equivalent Dataverse Web API v9.2 calls (representative routes)",
        D365_PROD_SAMPLES, cursor, total,
        intro="Base: https://{org}.api.crm.dynamics.com/api/data/v9.2/")

    wd_chunks = [WD_SAMPLES[i:i+4] for i in range(0, len(WD_SAMPLES), 4)]
    for i, chunk in enumerate(wd_chunks, start=1):
        slide_sample_cards(
            prs, "Workday · Samples",
            f"Internal endpoint request + response ({i}/{len(wd_chunks)})",
            chunk, cursor, total)
        cursor += 1

    cursor += slide_prod_samples(
        prs, "Workday · Production Samples",
        "Equivalent Workday REST + SOAP calls (representative routes)",
        WD_PROD_SAMPLES, cursor, total,
        intro="REST base: /ccx/api/{service}/v{n}/{tenant}/  ·  "
              "SOAP: /ccx/service/{tenant}/{service}/v{n}")

    actual_count = len(prs.slides)
    if actual_count != total:
        raise SystemExit(f"Slide count mismatch: built {actual_count}, expected {total}")
    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
