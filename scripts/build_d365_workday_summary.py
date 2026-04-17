"""Build exports/Integrations-D365-Workday-Summary.pptx — a 2-slide condensed
summary of the Dynamics 365 and Workday integrations (one slide per system).

Content is condensed from docs/integrations/api-overview.md and reuses the
visual language of scripts/build_integrations_deck.py. Run:
  python scripts/build_d365_workday_summary.py
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "exports" / "Integrations-D365-Workday-Summary.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

NAVY = RGBColor(0x12, 0x2B, 0x4A)
ACCENT = RGBColor(0xC2, 0x8A, 0x2A)
INK = RGBColor(0x1F, 0x24, 0x2E)
MUTED = RGBColor(0x5C, 0x66, 0x73)
LIGHT = RGBColor(0xF5, 0xF1, 0xE8)
LINE = RGBColor(0xD9, 0xD2, 0xC2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x1E, 0x26, 0x33)
CODE_FG = RGBColor(0xEC, 0xE6, 0xD0)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def text_box(slide, x, y, w, h, text, *, size=12, bold=False, color=INK,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        run.font.name = font; run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color
    return tb


def header(slide, eyebrow, title, idx, total):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
    fill(bar, ACCENT)
    text_box(slide, Inches(0.5), Inches(0.28), Inches(11), Inches(0.3),
             eyebrow.upper(), size=10, bold=True, color=ACCENT)
    text_box(slide, Inches(0.5), Inches(0.52), Inches(12.3), Inches(0.6),
             title, size=22, bold=True, color=NAVY)
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0.5), Inches(1.15), Inches(1.2), Emu(20000))
    fill(underline, NAVY)
    text_box(slide, Inches(11.7), Inches(7.1), Inches(1.5), Inches(0.3),
             f"{idx} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    text_box(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
             "DealPad · D365 + Workday Summary · April 2026", size=9, color=MUTED)


def section_label(slide, x, y, w, text):
    text_box(slide, x, y, w, Inches(0.28),
             text.upper(), size=10, bold=True, color=ACCENT)


def bullets(slide, x, y, w, h, items, *, size=11, color=INK, line_gap=3):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(line_gap)
        run = p.add_run(); run.text = f"•  {it}"
        run.font.name = "Calibri"; run.font.size = Pt(size); run.font.color.rgb = color


def code_block(slide, x, y, w, h, text, *, label=None):
    if label:
        text_box(slide, x, y - Inches(0.26), w, Inches(0.26),
                 label, size=9, bold=True, color=NAVY)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = NAVY
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.06)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run(); run.text = line if line else " "
        run.font.name = "Consolas"; run.font.size = Pt(8.5)
        run.font.color.rgb = CODE_FG


def api_table(slide, x, y, w, rows):
    h = Inches(0.22) * len(rows)
    tbl = slide.shapes.add_table(len(rows), 2, x, y, w, h).table
    tbl.columns[0].width = Inches(3.6)
    tbl.columns[1].width = w - Inches(3.6)
    for r, (left, right) in enumerate(rows):
        for c, val in enumerate((left, right)):
            cell = tbl.cell(r, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r % 2 == 0 else LIGHT
            cell.text = ""
            cell.margin_left = cell.margin_right = Inches(0.05)
            cell.margin_top = cell.margin_bottom = Inches(0.01)
            p = cell.text_frame.paragraphs[0]
            run = p.add_run(); run.text = val
            run.font.name = "Consolas"; run.font.size = Pt(8.5); run.font.color.rgb = INK


# ---------- slide builders ----------

def slide_d365(prs, idx, total):
    s = add_blank(prs)
    header(s, "Integration · Microsoft Dynamics 365",
           "CRM system of record for accounts and the opportunity pipeline",
           idx, total)

    # Context
    section_label(s, Inches(0.5), Inches(1.35), Inches(12.3), "Context")
    bullets(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.9), [
        "Inbound: pulls accounts and opportunities so DealPad can scope from real pipeline.",
        "Outbound: pushes fee, stage, probability and forecast back on every deal change.",
    ], size=11)

    # APIs (left column)
    section_label(s, Inches(0.5), Inches(2.55), Inches(6.0), "APIs · DealPad → Dataverse Web API v9.2")
    api_table(s, Inches(0.5), Inches(2.8), Inches(6.0), [
        ("GET  /api/dynamics/opportunities",  "GET /opportunities"),
        ("POST /api/dynamics/opportunities",  "POST /opportunities"),
        ("POST /api/dynamics/deals/:dealId/push", "PATCH /opportunities({opportunityid})"),
        ("POST /api/dynamics/sync",           "Bulk GET/PATCH (use $batch)"),
        ("GET  /api/dynamics/pipeline",       "GET /opportunities ?$select=…"),
    ])

    # Architecture (left column)
    section_label(s, Inches(0.5), Inches(4.15), Inches(6.0), "Architecture")
    bullets(s, Inches(0.5), Inches(4.4), Inches(6.0), Inches(2.5), [
        "Provider pattern: SimulatedDynamicsProvider today → LiveDynamicsProvider at cutover, no route changes.",
        "Triggers: outbound auto-push on stage / fee / margin change; inbound nightly batch + on-demand pull.",
        "Auth: OAuth 2.0 client-credentials via Azure AD (token → login.microsoftonline.com), audited in dynamics_sync_log.",
    ], size=10, line_gap=4)

    # Sample request / response (right column)
    section_label(s, Inches(6.85), Inches(2.55), Inches(6.0), "Sample · Create opportunity")
    code_block(s, Inches(6.85), Inches(2.85), Inches(6.0), Inches(1.85),
               "POST /api/dynamics/opportunities\n"
               "Content-Type: application/json\n\n"
               "{\n"
               '  "accountId": 42,\n'
               '  "name": "Crestwood Holdings - 2026 Audit",\n'
               '  "estimatedValue": 412000,\n'
               '  "stage": "Qualify"\n'
               "}",
               label="Request — DealPad")
    code_block(s, Inches(6.85), Inches(5.0), Inches(6.0), Inches(1.95),
               "{\n"
               '  "id": 137,\n'
               '  "opportunityNumber": "OPP-100204",\n'
               '  "estimatedValue": 412000,\n'
               '  "stage": "Qualify",\n'
               '  "probability": 20,\n'
               '  "forecastCategory": "Pipeline",\n'
               '  "syncStatus": "queued"\n'
               "}",
               label="Response")


def slide_workday(prs, idx, total):
    s = add_blank(prs)
    header(s, "Integration · Workday",
           "Source of truth for budgets, worker availability and standard cost rates",
           idx, total)

    section_label(s, Inches(0.5), Inches(1.35), Inches(12.3), "Context")
    bullets(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.9), [
        "Validates every deal save against cost-center budget, worker capacity and rate-card variance.",
        "Submission gate: over_budget or staffing_shortfall blocks unless Finance / Service Line Lead overrides with justification.",
    ], size=11)

    section_label(s, Inches(0.5), Inches(2.55), Inches(6.0), "APIs · DealPad → Workday REST + SOAP")
    api_table(s, Inches(0.5), Inches(2.8), Inches(6.0), [
        ("GET  /api/workday/cost-centers",       "GET /financialManagement/v1/{t}/costCenters"),
        ("GET  /api/workday/workers",            "GET /staffing/v6/{t}/workers"),
        ("PATCH /api/workday/rate-card/:id",     "SOAP Put_Compensation_Plan"),
        ("POST /api/workday/deals/:id/validate", "Composite: costCenters + workers + rules"),
        ("POST /api/workday/validations/:id/override", "DealPad audit (override fields)"),
    ])

    section_label(s, Inches(0.5), Inches(4.15), Inches(6.0), "Architecture")
    bullets(s, Inches(0.5), Inches(4.4), Inches(6.0), Inches(2.5), [
        "Provider pattern: SimulatedWorkdayProvider today → LiveWorkdayProvider at cutover, no route changes.",
        "Triggers: auto-validate on save, gate at submit, push approved project (committed budget reservation).",
        "Auth: OAuth 2.0 ISU for REST (or Basic auth for legacy SOAP); audited in workday_events + workday_validations.",
    ], size=10, line_gap=4)

    section_label(s, Inches(6.85), Inches(2.55), Inches(6.0), "Sample · Validate deal")
    code_block(s, Inches(6.85), Inches(2.85), Inches(6.0), Inches(1.85),
               "POST /api/workday/deals/87/validate\n"
               "Content-Type: application/json\n\n"
               "{ \"userName\": \"Sarah Chen\" }",
               label="Request — DealPad")
    code_block(s, Inches(6.85), Inches(5.0), Inches(6.0), Inches(1.95),
               "{\n"
               '  "ok": false,\n'
               '  "status": "staffing_shortfall",\n'
               '  "validationId": 412,\n'
               '  "summary": "Staffing shortfall: 240h.",\n'
               '  "findings": [\n'
               '    { "findingType": "staffing", "severity": "blocker",\n'
               '      "roleName": "Senior Consultant", "shortfallHours": 240 }\n'
               '  ]\n'
               "}",
               label="Response")


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    total = 2
    slide_d365(prs, 1, total)
    slide_workday(prs, 2, total)
    prs.save(OUT)
    print(f"Wrote {OUT.relative_to(ROOT)}  ({OUT.stat().st_size:,} bytes, {total} slides)")


if __name__ == "__main__":
    build()
