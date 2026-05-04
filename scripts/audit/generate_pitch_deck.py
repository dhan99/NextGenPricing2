#!/usr/bin/env python3
"""Generate docs/DealPad_Stakeholder_Pitch.pptx — phased delivery
narrative for the Armanino DealPad team.

Layout philosophy:
  - Tax-first throughout (Tax PDL is the primary persona; Tax PHB
    Standard Bundle, multi-entity 1040/1120/1065/1120S, batch
    renewals are the demo story).
  - AI introduced Phase 1 (heuristic, simulated narratives), real
    LLMs from Phase 2 onward, vector intelligence in Phase 3,
    trained ML in Phase 5.
  - Intapp lands in Phase 4 — Screening + Intake ONLY (DealPad
    remains the pricing engine; Intapp is the conflicts/intake
    source-of-truth). Bumped from earlier draft of Phase 5 so
    Armanino's team can lock workflow inside the project's runway.

Re-run from repo root:
    python3 scripts/audit/generate_pitch_deck.py
"""
from __future__ import annotations

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy

OUT_PATH = "docs/DealPad_Stakeholder_Pitch.pptx"

# ---- Brand palette (Armanino) ----
ARM_AMBER = RGBColor(0xDA, 0x72, 0x0F)        # primary
ARM_OLIVE = RGBColor(0x94, 0x93, 0x00)        # secondary
ARM_DARK = RGBColor(0x1F, 0x2A, 0x37)         # body text
ARM_MUTED = RGBColor(0x6B, 0x7A, 0x90)        # subdued text
ARM_LIGHT_AMBER = RGBColor(0xFF, 0xF1, 0xDA)  # background tint
ARM_LIGHT_OLIVE = RGBColor(0xF1, 0xF1, 0xDC)
ARM_BG = RGBColor(0xFA, 0xFA, 0xF7)            # off-white slide bg
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PHASE_GREEN = RGBColor(0x6B, 0x9E, 0x4F)
PHASE_BLUE = RGBColor(0x3B, 0x6E, 0xA8)
PHASE_PURPLE = RGBColor(0x7B, 0x4F, 0x9E)
PHASE_RED = RGBColor(0xC4, 0x6B, 0x4F)
PHASE_TEAL = RGBColor(0x4F, 0x9E, 0x9C)

# 16:9 widescreen
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


# ---------- helpers ----------

def add_blank_slide(prs):
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    # Background fill
    bg_fill = slide.background.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = ARM_BG
    return slide


def add_textbox(slide, x, y, w, h, text, *,
                font_size=18, bold=False, color=ARM_DARK,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bulletbox(slide, x, y, w, h, bullets, *,
                  font_size=14, color=ARM_DARK, bold_first=False,
                  bullet_color=None, indent_first_bold_line=True):
    """bullets: list of strings OR list of tuples (text, is_header)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, is_header = item
        else:
            text, is_header = item, False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        run = p.add_run()
        if is_header:
            run.text = text
            run.font.bold = True
            run.font.color.rgb = bullet_color or ARM_AMBER
            run.font.size = Pt(font_size + 1)
        else:
            run.text = "• " + text
            run.font.color.rgb = color
            run.font.size = Pt(font_size)
        run.font.name = "Calibri"
    return tb


def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_rounded_rect(slide, x, y, w, h, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.10
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_pill(slide, x, y, w, h, label, fill_color, text_color=WHITE, font_size=10):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = ""
    r = p.add_run()
    r.text = label
    r.font.name = "Calibri"
    r.font.size = Pt(font_size)
    r.font.bold = True
    r.font.color.rgb = text_color
    return shp


def add_footer(slide, idx, total, label="DealPad — Pricing & Scoping 2.0"):
    add_textbox(slide, Inches(0.4), Inches(7.10), Inches(10.0), Inches(0.3),
                label, font_size=9, color=ARM_MUTED)
    add_textbox(slide, Inches(12.4), Inches(7.10), Inches(0.7), Inches(0.3),
                f"{idx} / {total}", font_size=9, color=ARM_MUTED, align=PP_ALIGN.RIGHT)


def add_phase_band(slide, color, label):
    """Top band identifying which phase the slide belongs to."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.45), color)
    add_textbox(slide, Inches(0.4), Inches(0.05), Inches(8), Inches(0.35),
                label, font_size=14, bold=True, color=WHITE)


def add_h1(slide, text, y=Inches(0.55), color=ARM_DARK, font_size=30):
    add_textbox(slide, Inches(0.5), y, Inches(12.3), Inches(0.7),
                text, font_size=font_size, bold=True, color=color)


def add_h2(slide, text, y=Inches(1.12), color=ARM_OLIVE, font_size=15):
    add_textbox(slide, Inches(0.5), y, Inches(12.3), Inches(0.45),
                text, font_size=font_size, bold=False, color=color)


# ---------- slides ----------

def slide_cover(prs, total):
    slide = add_blank_slide(prs)
    # Hero band
    add_rect(slide, 0, 0, SLIDE_W, Inches(4.2), ARM_AMBER)
    add_rect(slide, 0, Inches(4.2), SLIDE_W, Inches(0.18), ARM_OLIVE)

    add_textbox(slide, Inches(0.6), Inches(0.6), Inches(12), Inches(0.5),
                "DealPad", font_size=20, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.6), Inches(1.15), Inches(12), Inches(1.2),
                "Pricing & Scoping 2.0", font_size=54, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.6), Inches(2.4), Inches(12), Inches(0.7),
                "From Excel sprawl to AI-native platform.",
                font_size=24, bold=False, color=WHITE)
    add_textbox(slide, Inches(0.6), Inches(2.95), Inches(12), Inches(0.6),
                "Tax-first.   Phased.   Indispensable.",
                font_size=22, bold=True, color=WHITE)

    # Sub-band
    add_textbox(slide, Inches(0.6), Inches(4.7), Inches(12), Inches(0.6),
                "A 5-phase plan to replace the Excel pricing workbooks with a",
                font_size=18, color=ARM_DARK)
    add_textbox(slide, Inches(0.6), Inches(5.05), Inches(12), Inches(0.6),
                "governed, AI-native platform integrated with the Quote-to-Cash stack.",
                font_size=18, color=ARM_DARK)

    add_textbox(slide, Inches(0.6), Inches(6.0), Inches(12), Inches(0.5),
                "Stakeholder briefing  |  Armanino DealPad team", font_size=14, color=ARM_MUTED)
    add_textbox(slide, Inches(0.6), Inches(6.4), Inches(12), Inches(0.5),
                "Live demo: https://dealpad-demo.onrender.com",
                font_size=12, color=ARM_OLIVE, bold=True)


def slide_problem(prs, idx, total):
    slide = add_blank_slide(prs)
    add_h1(slide, "The pain we're solving")
    add_h2(slide, "Pricing today is an Excel sprawl. It doesn't scale, audit, or learn.")

    # 3 columns of pain
    col_w = Inches(3.95)
    col_h = Inches(4.6)
    cols = [
        ("Calculation risk", PHASE_RED, [
            "200+ Excel workbooks across BUs",
            "No calc audit; one bad cell = wrong fee",
            "Each renewal cycle re-derives logic by hand",
            "Tax: 4-entity engagements stitched across tabs",
        ]),
        ("Process friction", ARM_AMBER, [
            "Manual approvals over email",
            "No version control on workbooks",
            "Tax-season renewals = days of copy-paste",
            "Prior-year baseline not auto-loaded",
        ]),
        ("Lost intelligence", PHASE_PURPLE, [
            "PDL judgment locked in spreadsheets",
            "Cannot search 'similar deals' across firm",
            "Margin patterns unrecoverable post-engagement",
            "Excel can't host AI assistance",
        ]),
    ]
    for i, (title, color, bullets) in enumerate(cols):
        x = Inches(0.5 + i * 4.2)
        y = Inches(1.7)
        add_rounded_rect(slide, x, y, col_w, col_h, ARM_LIGHT_AMBER if i == 1 else WHITE, line_color=color)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.15), col_w - Inches(0.3), Inches(0.5),
                    title, font_size=18, bold=True, color=color)
        add_bulletbox(slide, x + Inches(0.2), y + Inches(0.7),
                      col_w - Inches(0.3), col_h - Inches(0.8),
                      bullets, font_size=13)

    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.6),
                "The cost: margin leakage, audit exposure, and a ceiling on how much business Armanino can scoop up.",
                font_size=14, bold=True, color=ARM_DARK, align=PP_ALIGN.CENTER)
    add_footer(slide, idx, total)


def slide_opportunity(prs, idx, total):
    slide = add_blank_slide(prs)
    add_h1(slide, "The strategic opportunity")
    add_h2(slide, "One platform. Tax-first. AI-native. Quote-to-Cash integrated.")

    pillars = [
        ("Calc parity", "Reproduces every BU's Excel workbook outcomes,\ngolden-tested at the cent.", PHASE_GREEN),
        ("Tax-first design", "Multi-entity worksheets, assembly engine,\nbatch-renewal pipeline shipping in Phase 1.", ARM_AMBER),
        ("AI from day one", "Heuristic Phase 1 → real LLMs Phase 2 →\nvector intelligence Phase 3 → trained ML Phase 5.", PHASE_BLUE),
        ("Q2C integrations", "D365 (Phase 1) → Workday (Phase 3) →\nIntapp Intake + Screening (Phase 4) → Conga.", PHASE_PURPLE),
    ]
    for i, (title, body, color) in enumerate(pillars):
        x = Inches(0.5 + (i % 2) * 6.4)
        y = Inches(1.8 + (i // 2) * 2.4)
        add_rounded_rect(slide, x, y, Inches(6.0), Inches(2.0), WHITE, line_color=color)
        add_rect(slide, x, y, Inches(0.18), Inches(2.0), color)
        add_textbox(slide, x + Inches(0.4), y + Inches(0.2), Inches(5.5), Inches(0.5),
                    title, font_size=20, bold=True, color=color)
        add_textbox(slide, x + Inches(0.4), y + Inches(0.8), Inches(5.5), Inches(1.1),
                    body, font_size=14, color=ARM_DARK)

    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                "Outcome: governed, scaleable, AI-driven, audit-ready, indispensable.",
                font_size=13, bold=True, color=ARM_OLIVE, align=PP_ALIGN.CENTER)
    add_footer(slide, idx, total)


def slide_architecture(prs, idx, total):
    slide = add_blank_slide(prs)
    add_h1(slide, "Solution architecture at a glance")
    add_h2(slide, "Personas drive the workflow; integrations land in phases; AI sits horizontally.")

    # Three layers
    layers = [
        ("Personas", Inches(0.5), Inches(1.7), [
            "Project Delivery Lead (PDL)", "Service Line Lead (SLL)",
            "Pricing Operations (PO)", "Finance / FP&A",
            "Risk / QRM", "IT / Data Consumers",
        ], ARM_AMBER),
        ("DealPad core", Inches(4.5), Inches(1.7), [
            "Wizard (Setup → Approve)",
            "Multi-entity scope + assemblies",
            "Pricing engine + alt fee arrangements",
            "Approvals + audit + outbox",
            "AI engine (heuristic / LLM / vector / ML)",
        ], ARM_OLIVE),
        ("External systems", Inches(8.5), Inches(1.7), [
            "Dynamics 365  (Phase 1+)",
            "Workday Adaptive  (Phase 3+)",
            "Intapp Intake + Screening  (Phase 4)",
            "Conga CLM  (Phase 3+)",
            "Power BI  (Phase 4+)",
        ], PHASE_BLUE),
    ]
    for title, x, y, items, color in layers:
        add_rounded_rect(slide, x, y, Inches(3.8), Inches(4.5), WHITE, line_color=color)
        add_rect(slide, x, y, Inches(3.8), Inches(0.55), color)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.08), Inches(3.5), Inches(0.45),
                    title, font_size=17, bold=True, color=WHITE)
        add_bulletbox(slide, x + Inches(0.2), y + Inches(0.7),
                      Inches(3.5), Inches(3.7), items, font_size=13)

    # Bottom bar
    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
                "Outputs: branded proposals, engagement letters, Workday hours, Power BI feeds, AI cost telemetry.",
                font_size=12, color=ARM_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.4),
                "Foundation in place today (POC at https://dealpad-demo.onrender.com): every layer has been validated end-to-end.",
                font_size=12, bold=True, color=ARM_OLIVE, align=PP_ALIGN.CENTER)
    add_footer(slide, idx, total)


def slide_phase(prs, idx, total, *, num, color, title, weeks, headline,
                rows, demo, marquee_pillars):
    """
    rows: list of (lane, bullets) — typically 4 lanes (Pricing, AI, Integrations, Tax focus)
    demo: marquee demo description string
    marquee_pillars: 3 short phrases summarising why this phase wins
    """
    slide = add_blank_slide(prs)
    add_phase_band(slide, color, f"PHASE {num}  •  {weeks}")
    add_h1(slide, title, y=Inches(0.6))
    add_textbox(slide, Inches(0.5), Inches(1.18), Inches(12.3), Inches(0.5),
                headline, font_size=16, color=color, bold=True)

    # 4 lanes (2x2)
    lane_w = Inches(6.0)
    lane_h = Inches(2.05)
    for i, (lane_title, bullets) in enumerate(rows[:4]):
        x = Inches(0.5 + (i % 2) * 6.4)
        y = Inches(1.85 + (i // 2) * 2.20)
        add_rounded_rect(slide, x, y, lane_w, lane_h, WHITE, line_color=color)
        add_rect(slide, x, y, Inches(0.16), lane_h, color)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.08), lane_w - Inches(0.4), Inches(0.4),
                    lane_title, font_size=14, bold=True, color=color)
        add_bulletbox(slide, x + Inches(0.3), y + Inches(0.5),
                      lane_w - Inches(0.4), lane_h - Inches(0.55),
                      bullets, font_size=12)

    # Marquee demo strip
    yy = Inches(6.30)
    add_rect(slide, Inches(0.5), yy, Inches(12.3), Inches(0.65), ARM_LIGHT_AMBER)
    add_textbox(slide, Inches(0.65), yy + Inches(0.05), Inches(2.0), Inches(0.55),
                "Marquee demo →", font_size=12, bold=True, color=ARM_AMBER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(2.7), yy + Inches(0.05), Inches(10.0), Inches(0.55),
                demo, font_size=12, color=ARM_DARK, anchor=MSO_ANCHOR.MIDDLE)

    # Pillar pills
    pill_y = Inches(7.05)
    px = 0.5
    for p_label in marquee_pillars:
        add_pill(slide, Inches(px), pill_y, Inches(3.95), Inches(0.32), p_label, color, font_size=10)
        px += 4.10

    add_footer(slide, idx, total)


def slide_ai_layering(prs, idx, total):
    slide = add_blank_slide(prs)
    add_h1(slide, "AI layered phase-by-phase")
    add_h2(slide, "From heuristic seams to real LLMs to vectors to trained ML — without ripping out the stack.")

    phases = [
        ("Phase 1", "Heuristic AI",
         "Deal Similarity, Effort Estimation,\nMargin Advisor, Risk Summary —\nsimulated narratives, deterministic.",
         "Seams in place; no API key required.", PHASE_GREEN),
        ("Phase 2", "Real LLMs",
         "Anthropic Claude / OpenAI / Azure\npowering risk-summary, margin-advisor,\nscenario-recommendation narratives.",
         "AI telemetry dashboard live\nfrom day 1 of LLM wiring.", ARM_AMBER),
        ("Phase 3", "Vector + Voice + Predictive",
         "pgvector deal similarity (1536-dim,\n<500ms k-NN); voice-to-scope;\nscope-creep detector across 5 signals.",
         "Provider-swappable embeddings;\nML-model-ready evaluate() seam.", PHASE_BLUE),
        ("Phase 4", "Intapp AI integration",
         "Intake AI confidence-routed extraction:\nhigh-conf auto-progress,\nlow-conf reviewer matrix.",
         "AI sits ON the integration —\nIntake feeds DealPad scope.", PHASE_PURPLE),
        ("Phase 5", "Trained ML",
         "Effort estimator + margin LP solver\ntrained on Armanino's historical\nTax / Audit engagement history.",
         "Per-tenant tuning;\nmodel-version tracking.", PHASE_TEAL),
    ]
    box_w = Inches(2.41)
    box_h = Inches(4.7)
    gap = Inches(0.05)
    x0 = Inches(0.55)
    y0 = Inches(1.85)
    for i, (phase, headline, body, footnote, color) in enumerate(phases):
        x = x0 + (box_w + gap) * i
        # Card
        add_rounded_rect(slide, x, y0, box_w, box_h, WHITE, line_color=color)
        # Phase header
        add_rect(slide, x, y0, box_w, Inches(0.55), color)
        add_textbox(slide, x + Inches(0.1), y0 + Inches(0.08), box_w - Inches(0.2), Inches(0.45),
                    phase, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Headline
        add_textbox(slide, x + Inches(0.15), y0 + Inches(0.7),
                    box_w - Inches(0.3), Inches(0.7),
                    headline, font_size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        # Body
        add_textbox(slide, x + Inches(0.15), y0 + Inches(1.5),
                    box_w - Inches(0.3), Inches(2.3),
                    body, font_size=11, color=ARM_DARK)
        # Footnote
        add_rect(slide, x + Inches(0.15), y0 + Inches(3.85),
                 box_w - Inches(0.3), Inches(0.04), color)
        add_textbox(slide, x + Inches(0.15), y0 + Inches(3.95),
                    box_w - Inches(0.3), Inches(0.7),
                    footnote, font_size=10, color=ARM_MUTED)

    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
                "Same IntelligenceEngine + llm.ts seams from Phase 1; each phase swaps the implementation, not the contract.",
                font_size=12, bold=True, color=ARM_OLIVE, align=PP_ALIGN.CENTER)
    add_footer(slide, idx, total)


def slide_tax_thread(prs, idx, total):
    slide = add_blank_slide(prs)
    add_h1(slide, "Tax-first: how Tax differentiation grows phase by phase")
    add_h2(slide, "Tax is the demo headline at every milestone.")

    rows = [
        ("Phase 1", "Tax foundation", PHASE_GREEN, [
            "Tax scope catalog seeded:\n   1040, 1120, 1065, 1120S, Schedule K, multi-state",
            "Tax PHB Standard Bundle assembly with tier overrides",
            "Multi-entity worksheets — 4-entity Tax engagement on one deal",
            "Tax-specific prompts (entity type, complexity, prior-year)",
        ]),
        ("Phase 2", "Tax-season scale", ARM_AMBER, [
            "Batch renewal processor — 200+ Tax-renewal deals in one run",
            "Tax-tuned alt fee arrangements (compliance fixed-fee, retainer)",
            "Real LLM Risk Summary on Tax deals (Claude-powered narrative)",
            "Renewal Leadsheet with prior-year Tax baseline + YoY delta",
        ]),
        ("Phase 3", "Tax intelligence", PHASE_BLUE, [
            "Vector similarity tuned on Tax engagement embeddings",
            "Voice-to-scope from Tax client intake calls",
            "Scope-creep detector on Tax-renewal patterns",
            "Tax practice utilization → dynamic rate optimizer",
        ]),
        ("Phase 4", "Tax client lifecycle", PHASE_PURPLE, [
            "Tax client intake via Intapp → DealPad scope candidates",
            "Conflict screening on Tax engagements (pre-submit gate)",
            "Real-time collaboration: PDL + SLL on a Tax deal",
            "Client portal: Tax client reviews proposal via magic-link",
        ]),
        ("Phase 5", "Tax ML at production", PHASE_TEAL, [
            "Effort estimator trained on Armanino Tax history (<10% MAPE goal)",
            "Margin LP optimizer tuned on Tax practice constraints",
            "Per-tenant Tax catalog (multi-firm scaling)",
            "Slack/Teams alerts on Tax-deal approvals",
        ]),
    ]
    row_y = Inches(1.7)
    row_h = Inches(1.0)
    for i, (phase, title, color, items) in enumerate(rows):
        y = row_y + row_h * i
        add_rect(slide, Inches(0.5), y, Inches(1.5), row_h - Inches(0.05), color)
        add_textbox(slide, Inches(0.55), y + Inches(0.08), Inches(1.4), Inches(0.4),
                    phase, font_size=12, bold=True, color=WHITE)
        add_textbox(slide, Inches(0.55), y + Inches(0.45), Inches(1.4), Inches(0.5),
                    title, font_size=12, bold=True, color=WHITE)
        # Inline bullets, 2 per row to fit horizontally
        text = "  •  ".join(items[:4])
        # Actually use the bulletbox vertically
        add_rect(slide, Inches(2.1), y, Inches(10.7), row_h - Inches(0.05), WHITE, line_color=color)
        add_bulletbox(slide, Inches(2.2), y + Inches(0.05),
                      Inches(10.5), row_h - Inches(0.1), items, font_size=10)

    add_footer(slide, idx, total)


def slide_alternatives(prs, idx, total):
    slide = add_blank_slide(prs)
    add_h1(slide, "DealPad vs the alternatives")
    add_h2(slide, "Why building forward is the only path that compounds.")

    cols = [
        ("Status quo (Excel)", PHASE_RED, [
            ("✘", "No governance, no audit"),
            ("✘", "Calc errors per workbook"),
            ("✘", "Manual renewals at scale"),
            ("✘", "Cannot host AI"),
            ("✘", "Tax 4-entity = tab sprawl"),
            ("✘", "Lost institutional knowledge"),
        ]),
        ("Intapp Pricing engine", PHASE_BLUE, [
            ("◐", "Vendor pricing logic — not Armanino's Excel parity"),
            ("◐", "Generic; no Tax-first multi-entity"),
            ("◐", "AI = vendor roadmap, not yours"),
            ("✘", "Doesn't replace Excel logic 1:1"),
            ("✓", "Stays as Intake + Screening source-of-truth (DealPad integrates)"),
            ("✘", "Build-vs-buy lock-in"),
        ]),
        ("DealPad", PHASE_GREEN, [
            ("✓", "Calc parity to your workbooks"),
            ("✓", "Tax-first multi-entity + assemblies"),
            ("✓", "AI Phase 1 → trained ML Phase 5"),
            ("✓", "Replaces Excel, integrates Intapp Intake/Screening"),
            ("✓", "Already 80% built (POC live today)"),
            ("✓", "Built around your personas"),
        ]),
    ]
    for i, (title, color, items) in enumerate(cols):
        x = Inches(0.5 + i * 4.2)
        y = Inches(1.7)
        h = Inches(5.0)
        add_rounded_rect(slide, x, y, Inches(4.0), h, WHITE, line_color=color)
        add_rect(slide, x, y, Inches(4.0), Inches(0.55), color)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.08), Inches(3.7), Inches(0.5),
                    title, font_size=15, bold=True, color=WHITE)
        # items
        cy = y + Inches(0.7)
        for sym, text in items:
            tb = slide.shapes.add_textbox(x + Inches(0.2), cy, Inches(3.7), Inches(0.7))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
            tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r1 = p.add_run()
            r1.text = sym + "  "
            r1.font.bold = True
            r1.font.size = Pt(13)
            sym_color = PHASE_GREEN if sym == "✓" else (PHASE_RED if sym == "✘" else ARM_AMBER)
            r1.font.color.rgb = sym_color
            r2 = p.add_run()
            r2.text = text
            r2.font.size = Pt(12)
            r2.font.color.rgb = ARM_DARK
            cy += Inches(0.62)

    add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                "DealPad complements Intapp — pricing engine in DealPad, intake & screening in Intapp. Not a competition.",
                font_size=13, bold=True, color=ARM_OLIVE, align=PP_ALIGN.CENTER)
    add_footer(slide, idx, total)


def slide_investment(prs, idx, total):
    slide = add_blank_slide(prs)
    add_h1(slide, "Investment summary")
    add_h2(slide, "5 phases × 8 weeks = 40 weeks. Milestone-gated GO decisions at each handoff.")

    # Timeline
    timeline_y = Inches(1.85)
    bar_h = Inches(0.55)
    phases = [
        ("P1", "Excel Parity\n+ Tax Foundation", PHASE_GREEN),
        ("P2", "Pricing\nSophistication\n+ Real LLMs", ARM_AMBER),
        ("P3", "Intelligence\nEngine", PHASE_BLUE),
        ("P4", "Intapp +\nClient Maturity", PHASE_PURPLE),
        ("P5", "Production\n+ Trained ML", PHASE_TEAL),
    ]
    bar_w = Inches(2.5)
    bar_x = Inches(0.5)
    for i, (label, name, color) in enumerate(phases):
        x = bar_x + bar_w * i
        add_rect(slide, x, timeline_y, bar_w - Inches(0.05), bar_h, color)
        add_textbox(slide, x + Inches(0.1), timeline_y + Inches(0.08),
                    bar_w - Inches(0.25), Inches(0.4),
                    label + " · 8 wks", font_size=11, bold=True, color=WHITE)
        add_textbox(slide, x + Inches(0.1), timeline_y + bar_h + Inches(0.08),
                    bar_w - Inches(0.25), Inches(0.9), name,
                    font_size=10, color=ARM_DARK)

    # Team shape
    add_textbox(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.4),
                "Team shape", font_size=16, bold=True, color=ARM_OLIVE)
    team = [
        ("Product Manager", "Full-time, all phases"),
        ("Senior Engineers (2)", "Full-time, all phases"),
        ("AI/ML Engineer", "Phase 2 onward (full-time from P3)"),
        ("UX Designer", "Phase 1, 4 (UI-heavy work)"),
        ("Armanino SME (PDL+PO)", "Part-time embed, all phases — calc parity acceptance"),
    ]
    ty = Inches(4.2)
    for label, when in team:
        add_textbox(slide, Inches(0.7), ty, Inches(4), Inches(0.35),
                    "•  " + label, font_size=12, bold=True, color=ARM_DARK)
        add_textbox(slide, Inches(4.5), ty, Inches(8), Inches(0.35),
                    when, font_size=12, color=ARM_MUTED)
        ty += Inches(0.35)

    # ROI levers
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.4),
                "ROI levers (Armanino to validate magnitudes)",
                font_size=14, bold=True, color=ARM_OLIVE)
    roi = [
        "PDL time saved: 4–6 hrs / Tax engagement × 1,000+ engagements / yr",
        "Margin recovery: 1–2% lift × pipeline through advisor + creep detector",
        "Audit defensibility: zero-cost audit-trail vs ad-hoc Excel reconstruction",
        "Renewal throughput: 200-deal batch in hours, not days",
    ]
    add_bulletbox(slide, Inches(0.7), Inches(6.4), Inches(12), Inches(0.6),
                  roi, font_size=11)

    add_footer(slide, idx, total)


def slide_ask(prs, idx, total):
    slide = add_blank_slide(prs)
    add_h1(slide, "What we need from you — and what you get")
    add_h2(slide, "Phase 1 GO is a small, low-risk commitment with measurable proof in 8 weeks.")

    # 2-column layout
    left_x = Inches(0.5)
    right_x = Inches(6.9)
    box_w = Inches(6.0)
    box_h = Inches(4.5)
    y = Inches(1.85)
    add_rounded_rect(slide, left_x, y, box_w, box_h, ARM_LIGHT_AMBER, line_color=ARM_AMBER)
    add_textbox(slide, left_x + Inches(0.2), y + Inches(0.15), box_w - Inches(0.3), Inches(0.5),
                "What we need from Armanino", font_size=18, bold=True, color=ARM_AMBER)
    asks = [
        "10 representative Excel pricing workbooks for calc-parity",
        "Tax scope-catalog source data (1040 / 1120 / 1065 / 1120S\n   + multi-state list + Tax PHB Standard Bundle definition)",
        "D365 sandbox access for opportunity-context auto-load",
        "PDL + PO part-time embed for calc-parity acceptance",
        "GO decision at end of each 8-week phase (5 GO gates total)",
        "(Phase 2) Anthropic / OpenAI / Azure OpenAI key access",
        "(Phase 4) Intapp Intake + Screening API access",
        "(Phase 5) Historical engagement data for ML training",
    ]
    add_bulletbox(slide, left_x + Inches(0.2), y + Inches(0.7),
                  box_w - Inches(0.3), box_h - Inches(0.8), asks, font_size=12)

    add_rounded_rect(slide, right_x, y, box_w, box_h, ARM_LIGHT_OLIVE, line_color=ARM_OLIVE)
    add_textbox(slide, right_x + Inches(0.2), y + Inches(0.15), box_w - Inches(0.3), Inches(0.5),
                "What you get", font_size=18, bold=True, color=ARM_OLIVE)
    gets = [
        "Week 4 — Shared dev URL with Tax scope catalog seeded",
        "Week 8 — Calc parity sign-off; Phase 1 demo to leadership",
        "Week 16 — Real LLMs powering risk + margin AI",
        "Week 24 — Vector similarity + voice-to-scope live",
        "Week 32 — Intapp Intake/Screening live + client portal",
        "Week 40 — Trained ML + Slack/Teams + production hardened",
        "Always — open source code, audit trail, no vendor lock-in",
    ]
    add_bulletbox(slide, right_x + Inches(0.2), y + Inches(0.7),
                  box_w - Inches(0.3), box_h - Inches(0.8), gets, font_size=12)

    # Bottom CTA
    add_rect(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.65), ARM_AMBER)
    add_textbox(slide, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.55),
                "Ask: 8-week Phase 1 GO. Calc parity proven. Tax foundation seeded. AI-native demo by Week 9.",
                font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, idx, total)


# ---------- main ----------

def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    TOTAL = 14
    idx = 0

    # ---- Cover ----
    idx += 1; slide_cover(prs, TOTAL)
    # ---- Problem ----
    idx += 1; slide_problem(prs, idx, TOTAL)
    # ---- Opportunity ----
    idx += 1; slide_opportunity(prs, idx, TOTAL)
    # ---- Architecture ----
    idx += 1; slide_architecture(prs, idx, TOTAL)

    # ---- Phase 1 ----
    idx += 1
    slide_phase(prs, idx, TOTAL,
        num=1, color=PHASE_GREEN,
        title="Phase 1 — Excel Parity + Tax Foundation",
        weeks="8–10 weeks",
        headline="Same numbers as your Excel workbooks. Tax-first scope. Heuristic AI from day one.",
        rows=[
            ("Pricing engine + workflow", [
                "Calc-parity golden test pinned against current Excel",
                "8-step wizard (Setup → Approve), versioning, audit trail",
                "RBAC across 6 personas + tiered approvals",
                "T&M arrangement; engagement letter export (HTML/PDF)",
            ]),
            ("Tax foundation", [
                "Scope catalog: 1040 / 1120 / 1065 / 1120S / Schedule K",
                "Tax PHB Standard Bundle assembly + tier overrides",
                "Multi-entity worksheets (4 entity tabs on one deal)",
                "Tax-specific prompt set (entity type, complexity, PY ref)",
            ]),
            ("AI (heuristic)", [
                "Deal Similarity, Effort Estimation, Margin Advisor",
                "Risk Summary with simulated narrative",
                "All seams in place — no API key dependency yet",
                "AI telemetry table provisioned, ready for Phase 2",
            ]),
            ("Integrations", [
                "Dynamics 365 inbound: account / opportunity context",
                "Auto-create DealPad client + deal on import",
                "D365 owners + quotas synced for routing visibility",
                "Workday + Intapp + Conga: read-only catalog only",
            ]),
        ],
        demo="Import Crestwood Holdings opp from D365 → walk a 4-entity Tax PHB engagement end-to-end → prove calc parity against the equivalent Excel workbook on the same screen.",
        marquee_pillars=[
            "Calc parity = credibility",
            "Tax-first design",
            "AI seams ready",
        ],
    )

    # ---- Phase 2 ----
    idx += 1
    slide_phase(prs, idx, TOTAL,
        num=2, color=ARM_AMBER,
        title="Phase 2 — Pricing Sophistication + Real LLMs",
        weeks="8 weeks",
        headline="Tax-season scale. Six fee models. Real LLM narratives powering approval AI.",
        rows=[
            ("Pricing depth", [
                "Assembly engine (math.js sandbox; PHB → 87 line items / tier)",
                "All 6 fee arrangements: Fixed / Capped / Contingent /",
                "   Retainer / Hybrid / T&M",
                "Per-deal margin-target overrides + engagement-input policy",
            ]),
            ("Tax-season ops", [
                "Batch renewal processor — 200+ Tax deals in one job",
                "Adjustment rules (rate uplift, hour adjust, margin override)",
                "Renewal Leadsheet with prior-year baseline + YoY delta",
                "Margin-targets governance (firm/BU/serviceLine)",
            ]),
            ("AI (real LLMs)", [
                "Anthropic Claude / OpenAI / Azure flag-controlled",
                "Risk Summary, Margin Advisor, Scenario Recommendation",
                "   power structured narratives via llm.completeStructured",
                "AI telemetry dashboard live: token cost, p95 latency, errors",
            ]),
            ("Integrations", [
                "D365 bi-directional: status / owner / amount push back",
                "   on changes; auto-push on approval transitions",
                "Workday cost-center read for visibility",
                "Conga CLM template selection on engagement letter step",
            ]),
        ],
        demo="Run the Tax-season renewal batch over 50 Tax deals with a 5% rate uplift + tier override rule → flagged variances → real LLM-narrated risk summary on a flagged deal.",
        marquee_pillars=[
            "Tax-season scale",
            "Six fee models",
            "Live LLM AI",
        ],
    )

    # ---- Phase 3 ----
    idx += 1
    slide_phase(prs, idx, TOTAL,
        num=3, color=PHASE_BLUE,
        title="Phase 3 — Intelligence Engine: Vectors, Voice, Watchdogs",
        weeks="8 weeks",
        headline="Sub-500ms 'similar deals.' Voice memos become scope. Scope-creep flagged before deals miss numbers.",
        rows=[
            ("Vector intelligence", [
                "pgvector deal similarity (1536-dim embeddings)",
                "Sub-500ms k-NN against thousands of historical deals",
                "Lazy backfill on first query; provider-swappable embeddings",
                "Tax engagement embeddings for Tax-context similarity",
            ]),
            ("Predictive AI", [
                "Voice-to-scope: audio → transcription → scope candidates",
                "Predictive Scope Creep Detector: 5 signal types,",
                "   ML-swap-ready evaluate() seam",
                "Dynamic rate optimizer (capacity + velocity + margin)",
            ]),
            ("Time + budget", [
                "Time tracking + AI-suggested time entries",
                "Budget-to-actuals monitoring with thresholds + alerts",
                "Time-entry-driven actuals (replaces heuristic projection)",
                "Cost / fee variance tracked at deal + entity granularity",
            ]),
            ("Integrations", [
                "Workday push: project commitment + cost-center bump on",
                "   approval; idempotent via activity_log guard",
                "Conga CLM engagement-letter delivery (email/eSign/portal)",
                "Power BI feed (foundation only; full feed in Phase 4)",
            ]),
        ],
        demo="PDL records voice memo from Tax client meeting → scope items extracted with confidence scores → applied to deal in 90 seconds, with vector similarity surfacing 5 historical Tax deals at this margin profile.",
        marquee_pillars=[
            "Vector similarity",
            "Voice → scope",
            "Predictive watchdog",
        ],
    )

    # ---- Phase 4 (Intapp moved here) ----
    idx += 1
    slide_phase(prs, idx, TOTAL,
        num=4, color=PHASE_PURPLE,
        title="Phase 4 — Intapp Integration + Client Maturity",
        weeks="8 weeks",
        headline="Intapp Intake + Screening as primary touch points. Live client portal. Real-time PDL collaboration.",
        rows=[
            ("Intapp Intake", [
                "AI confidence-routed extraction:",
                "   high-conf intakes auto-progress; low-conf → reviewer",
                "Tax client intake → DealPad scope candidates auto-seeded",
                "Reviewer matrix mapped to PDL / SLL / QRM personas",
            ]),
            ("Intapp Risk Screening", [
                "Bi-directional: nightly re-screen + screening results sync",
                "Conflict + mitigation lifecycle integrated into approval",
                "Outcome push on approval/rejection; mitigation status",
                "Pre-submit gate blocks submission on conflicts (overridable)",
            ]),
            ("Client + collaboration", [
                "Magic-link client portal (SHA-256 hashed tokens; no account)",
                "Read-only proposal + scope review for clients",
                "Yjs real-time collaborative scoping (multi-PDL editing)",
                "Domain-driven refactor wraps Phase 4 (Service Bus / Kafka ready)",
            ]),
            ("Reporting + governance", [
                "Power BI data feeds for renewal tracking",
                "Configurable approval workflows per BU / exception type",
                "Out-of-scope add-on separation; multi-year / term modeling",
                "Geo / offshore discount engine; Tax-firm geo applied",
            ]),
        ],
        demo="New Tax client → Intapp Intake → DealPad deal seeded with scope → Intapp screens for conflicts → PDL + SLL collaborate live on scope → Tax client reviews proposal via magic-link.",
        marquee_pillars=[
            "Intapp Intake live",
            "Intapp Screening live",
            "Real-time collaboration",
        ],
    )

    # ---- Phase 5 ----
    idx += 1
    slide_phase(prs, idx, TOTAL,
        num=5, color=PHASE_TEAL,
        title="Phase 5 — Production Hardening + Trained ML",
        weeks="8 weeks",
        headline="Models trained on YOUR Tax history. Multi-region production. Slack/Teams. Ready for the long tail.",
        rows=[
            ("Trained ML", [
                "Effort estimator: sklearn / Azure ML on Armanino Tax history",
                "Margin LP optimizer replaces heuristic",
                "Per-tenant tuning; model-version tracking in AI telemetry",
                "Goal: <10% MAPE on effort estimation across Tax engagements",
            ]),
            ("Notifications", [
                "Slack native app: approval prompts + slash commands",
                "Teams native app: channel-bound notifications",
                "Scope-creep + budget alerts pushed to ops channels",
                "Batch-renewal completion notifications",
            ]),
            ("Production hardening", [
                "Multi-region failover; observability stack (Datadog or eq.)",
                "Per-deal AI cost attribution; model-version tracking",
                "SOC compliance posture; audit retention policy",
                "Configurable approval workflows + exception thresholds",
            ]),
            ("Long-tail features", [
                "Out-of-scope add-on tracking + approval acknowledgement",
                "Multi-year / term modeling for Tax retainers + price escalators",
                "Geo / offshore discount rules engine",
                "Fast-track approval + delegation / backup approvers",
            ]),
        ],
        demo="Trained effort estimator validated against 100 historical Tax engagements (<10% MAPE) → per-deal AI cost report → Slack approval flow → multi-region failover drill.",
        marquee_pillars=[
            "Tax-trained ML",
            "Slack / Teams",
            "Production-hardened",
        ],
    )

    # ---- AI layering visual ----
    idx += 1; slide_ai_layering(prs, idx, TOTAL)

    # ---- Tax thread ----
    idx += 1; slide_tax_thread(prs, idx, TOTAL)

    # ---- Alternatives ----
    idx += 1; slide_alternatives(prs, idx, TOTAL)

    # ---- Investment ----
    idx += 1; slide_investment(prs, idx, TOTAL)

    # ---- Ask ----
    idx += 1; slide_ask(prs, idx, TOTAL)

    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH} — {idx} slides")


if __name__ == "__main__":
    main()
