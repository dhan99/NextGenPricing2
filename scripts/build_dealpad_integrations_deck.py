"""
Build the DealPad Integrations Architecture deck — single condensed slide.

One executive slide that brings together all five integrations
(Dynamics 365, Workday, Conga, Intapp, Power BI) with:
- A compact DealPad hub banner at the top.
- Five side-by-side integration cards underneath, each showing:
  direction badge, capability hook, outbound + inbound payload + trigger,
  technology stack chips, and the business / sales-closure context.

Content sourced from client/src/pages/ArchitectureInteractive.tsx (the
in-app Architecture Hub) and server/routes.ts trigger points.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---------- palette ----------
NAVY = RGBColor(0x0F, 0x2A, 0x4A)
NAVY_DEEP = RGBColor(0x0A, 0x1F, 0x3A)
INK = RGBColor(0x1A, 0x2B, 0x3C)
SUB = RGBColor(0x55, 0x6B, 0x82)
MUTED = RGBColor(0x6B, 0x7B, 0x8C)
RULE = RGBColor(0xE3, 0xE8, 0xEF)
BG = RGBColor(0xF6, 0xF8, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1E, 0x8E, 0x5A)
AMBER = RGBColor(0xD8, 0x8A, 0x0E)

INTEGRATIONS = [
    {
        "key": "dynamics",
        "name": "Dynamics 365",
        "tag": "CRM • opportunity",
        "color": RGBColor(0x00, 0x78, 0xD4),
        "direction": "Bi-directional",
        "capability": "Opportunity sync + outcome push on every approval",
        "outbound": "Status, fee, cost, margin %, owner",
        "outbound_trigger": "autoPushDeal() on every approval transition",
        "inbound": "Accounts, opportunities, contacts, owner",
        "inbound_trigger": "On deal create + manual refresh",
        "tech": ["Dataverse Web API v9.2", "OAuth 2.0", "REST"],
        "business": "Every deal anchored to a Dynamics opportunity — opportunity stage and revenue auto-update on Approved / Lost / Withdrawn.",
    },
    {
        "key": "workday",
        "name": "Workday",
        "tag": "Budget • resourcing",
        "color": RGBColor(0x00, 0x86, 0xA8),
        "direction": "Bi-directional",
        "capability": "Validation gate + project create + budget reserve",
        "outbound": "Project record + committed-budget reserve",
        "outbound_trigger": "autoPushWorkdayProject() on Approved (transactional sentinel)",
        "inbound": "Cost centers, workers, rate cards by skill/grade",
        "inbound_trigger": "Nightly sync + on-demand during scoping",
        "tech": ["REST + SOAP", "FM • Staffing • Comp", "ISU / OAuth 2.0"],
        "business": "Rate cards drive the pricing grid; submission blocked on rate variance, budget headroom, or staffing shortfall (override needs Finance/SLL justification).",
    },
    {
        "key": "conga",
        "name": "Conga CLM",
        "tag": "Engagement letters",
        "color": RGBColor(0xC0, 0x39, 0x2B),
        "direction": "Bi-directional",
        "capability": "Template letters + e-sign delivery + status reconcile",
        "outbound": "Letter generation + delivery request",
        "outbound_trigger": "Auto on Approved + manual 'Generate Letter'",
        "inbound": "Delivery status (sent / signed / delivered)",
        "inbound_trigger": "Webhook + polling on outstanding letters",
        "tech": ["Conga REST API", "DocuSign / Adobe Sign", "Email / portal"],
        "business": "Sales closure: approved deal becomes a governed engagement letter in seconds; engagement_letters row flips to 'delivered' once client signs.",
    },
    {
        "key": "intapp",
        "name": "Intapp Risk",
        "tag": "Conflicts • independence",
        "color": RGBColor(0x6B, 0x3F, 0xA0),
        "direction": "Bi-directional",
        "capability": "Conflict / independence screening + outcome push",
        "outbound": "Mitigation decisions + final deal outcome",
        "outbound_trigger": "On resolve/waive/reject + autoPushIntappOutcome()",
        "inbound": "Screening results, mitigation hits, severity",
        "inbound_trigger": "On deal create / client change",
        "tech": ["Intapp Risk REST API", "Conflicts + Independence", "OAuth 2.0"],
        "business": "Risk gate: no approval while open hits remain. Engagement onboarding opens in Intapp the moment DealPad approves the deal.",
    },
    {
        "key": "powerbi",
        "name": "Power BI",
        "tag": "Analytics • benchmarks",
        "color": RGBColor(0xC9, 0xA2, 0x27),
        "direction": "One-way (read)",
        "capability": "Executive dashboards + margin benchmarks + DealPad-vs-Workday actuals",
        "outbound": "—  (Power BI is read-only)",
        "outbound_trigger": "—",
        "inbound": "Deals, scope, pricing, scenarios, approvals, sync events",
        "inbound_trigger": "Hourly DirectQuery + nightly scheduled refresh",
        "tech": ["Power BI Service", "DirectQuery + Refresh", "RLS via persona claims"],
        "business": "Pipeline & forecast, cycle-time analytics (opportunity → approved → signed), and the benchmark feed for the AI margin advisor.",
    },
]

# ---------- helpers ----------
def add_text(slide, x, y, w, h, text, *, size=12, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb

def add_rect(slide, x, y, w, h, *, fill=WHITE, line=RULE, line_w=0.75,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.05):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: s.adjustments[0] = adj
        except: pass
    s.text_frame.text = ""
    return s

def draw_arrow(slide, x1, y1, x2, y2, color, *, weight=1.5, dashed=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    ln = conn.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    ln.append(ln.makeelement(qn('a:tailEnd'),
              {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return conn

def label_row(slide, x, y, w, label, value, *, label_color, value_color=INK,
              label_size=8, value_size=9, line_h=Inches(0.36)):
    """Small two-line block: ALL-CAPS label on top, value beneath."""
    add_text(slide, x, y, w, Inches(0.16), label.upper(),
             size=label_size, bold=True, color=label_color)
    add_text(slide, x, y + Inches(0.16), w, line_h - Inches(0.16),
             value, size=value_size, color=value_color)

# ---------- presentation ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, SW, SH, fill=BG, line=BG, shape=MSO_SHAPE.RECTANGLE)

# --- Title ---
add_text(slide, Inches(0.5), Inches(0.28), SW - Inches(1.0), Inches(0.5),
         "DealPad Integrations Architecture", size=24, bold=True, color=NAVY)
add_text(slide, Inches(0.5), Inches(0.78), SW - Inches(1.0), Inches(0.32),
         "How DealPad connects to the Quote-to-Cash stack — bi-directional with every system except Power BI (read-only).",
         size=11, color=SUB, italic=True)

# Legend (top-right)
lg_y = Inches(0.32)
# solid swatch
sw1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
    SW - Inches(4.6), lg_y + Inches(0.12),
    SW - Inches(4.2), lg_y + Inches(0.12))
sw1.line.color.rgb = NAVY; sw1.line.width = Pt(1.75)
ln = sw1.line._get_or_add_ln()
ln.append(ln.makeelement(qn('a:tailEnd'),
          {'type': 'triangle', 'w': 'med', 'len': 'med'}))
add_text(slide, SW - Inches(4.15), lg_y + Inches(0.02), Inches(1.7), Inches(0.22),
         "Outbound from DealPad", size=8, color=INK)
sw2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
    SW - Inches(2.4), lg_y + Inches(0.12),
    SW - Inches(2.0), lg_y + Inches(0.12))
sw2.line.color.rgb = RGBColor(0x8A, 0x9B, 0xB0); sw2.line.width = Pt(1.25)
ln2 = sw2.line._get_or_add_ln()
ln2.append(ln2.makeelement(qn('a:prstDash'), {'val': 'dash'}))
ln2.append(ln2.makeelement(qn('a:tailEnd'),
          {'type': 'triangle', 'w': 'med', 'len': 'med'}))
add_text(slide, SW - Inches(1.95), lg_y + Inches(0.02), Inches(1.7), Inches(0.22),
         "Inbound to DealPad", size=8, color=INK)

# --- DealPad hub banner ---
hub_top = Inches(1.18)
hub_h = Inches(0.62)
hub = add_rect(slide, Inches(0.5), hub_top, SW - Inches(1.0), hub_h,
               fill=NAVY, line=NAVY_DEEP, line_w=1.0, adj=0.25)
tf = hub.text_frame
tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
tf.margin_top = 0; tf.margin_bottom = 0
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = "DealPad"
r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
r2 = p.add_run()
r2.text = "    Pricing & Scoping 2.0 hub  •  Express.js API  •  PostgreSQL  •  AI services  •  RBAC + audit log"
r2.font.size = Pt(11); r2.font.color.rgb = RGBColor(0xCB, 0xDA, 0xEE); r2.font.name = "Calibri"

# Connector lines from hub down to each card (visual hub-and-spoke)
hub_bottom_y = hub_top + hub_h
card_top = hub_bottom_y + Inches(0.45)

n = len(INTEGRATIONS)
side_margin = Inches(0.5)
gap = Inches(0.12)
total_w = SW - side_margin * 2
card_w = (total_w - gap * (n - 1)) / n
card_h = SH - card_top - Inches(0.55)

for i, integ in enumerate(INTEGRATIONS):
    x = side_margin + i * (card_w + gap)
    # spoke connector
    cx = int(x + card_w / 2)
    spoke = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
        cx, int(hub_bottom_y), cx, int(card_top))
    spoke.line.color.rgb = integ["color"]; spoke.line.width = Pt(1.5)

    # Card
    card = add_rect(slide, int(x), int(card_top), int(card_w), int(card_h),
                    fill=WHITE, line=RULE, line_w=0.75, adj=0.04)

    # Header strip
    hdr_h = Inches(0.62)
    hdr = add_rect(slide, int(x), int(card_top), int(card_w), int(hdr_h),
                   fill=integ["color"], line=integ["color"], line_w=0.75, adj=0.04)
    # Mask the bottom rounded corners by overlaying a thin rectangle (so the
    # header reads as a top band)
    mask = add_rect(slide,
        int(x), int(card_top + hdr_h - Inches(0.18)),
        int(card_w), int(Inches(0.18)),
        fill=integ["color"], line=integ["color"], line_w=0,
        shape=MSO_SHAPE.RECTANGLE)

    # Header text
    tf = hdr.text_frame
    tf.margin_left = Inches(0.16); tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.06); tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = integ["name"]
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(0)
    r2 = p2.add_run(); r2.text = integ["tag"]
    r2.font.size = Pt(9); r2.font.italic = True
    r2.font.color.rgb = RGBColor(0xF0, 0xF4, 0xFA); r2.font.name = "Calibri"

    # Direction badge (top-right of header)
    badge_w = Inches(1.05); badge_h = Inches(0.24)
    one_way = integ["direction"].lower().startswith("one")
    badge = add_rect(slide,
        int(x + card_w - badge_w - Inches(0.12)),
        int(card_top + Inches(0.12)),
        int(badge_w), int(badge_h),
        fill=WHITE, line=WHITE, line_w=0, adj=0.5)
    btf = badge.text_frame
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(btf, m, 0)
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = "BI-DIRECTIONAL" if not one_way else "ONE-WAY (READ)"
    br.font.size = Pt(7); br.font.bold = True
    br.font.color.rgb = integ["color"]; br.font.name = "Calibri"

    # ---- Body content ----
    pad = Inches(0.18)
    body_x = int(x + pad)
    body_w = int(card_w - pad * 2)
    cur_y = card_top + hdr_h + Inches(0.14)

    # Capability hook
    add_text(slide, body_x, int(cur_y), body_w, Inches(0.16), "CAPABILITY",
             size=8, bold=True, color=integ["color"])
    cur_y += Inches(0.16)
    cap_box = add_text(slide, body_x, int(cur_y), body_w, Inches(0.55),
                       integ["capability"], size=10, bold=True, color=INK)
    cur_y += Inches(0.58)

    # Outbound mini-row with arrow icon
    add_text(slide, body_x, int(cur_y), body_w, Inches(0.18),
             "→  OUTBOUND", size=8, bold=True, color=NAVY)
    cur_y += Inches(0.18)
    add_text(slide, body_x, int(cur_y), body_w, Inches(0.32),
             integ["outbound"], size=9, color=INK)
    cur_y += Inches(0.32)
    add_text(slide, body_x, int(cur_y), body_w, Inches(0.30),
             integ["outbound_trigger"], size=8, color=MUTED, italic=True)
    cur_y += Inches(0.32)

    # Inbound mini-row
    in_color = MUTED if one_way else integ["color"]
    add_text(slide, body_x, int(cur_y), body_w, Inches(0.18),
             "←  INBOUND", size=8, bold=True, color=in_color)
    cur_y += Inches(0.18)
    add_text(slide, body_x, int(cur_y), body_w, Inches(0.32),
             integ["inbound"], size=9, color=INK)
    cur_y += Inches(0.32)
    add_text(slide, body_x, int(cur_y), body_w, Inches(0.30),
             integ["inbound_trigger"], size=8, color=MUTED, italic=True)
    cur_y += Inches(0.36)

    # Divider
    div = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
        body_x, int(cur_y), body_x + body_w, int(cur_y))
    div.line.color.rgb = RULE; div.line.width = Pt(0.5)
    cur_y += Inches(0.10)

    # Tech stack
    add_text(slide, body_x, int(cur_y), body_w, Inches(0.18),
             "TECH STACK", size=8, bold=True, color=integ["color"])
    cur_y += Inches(0.18)
    tech_text = "  •  ".join(integ["tech"])
    add_text(slide, body_x, int(cur_y), body_w, Inches(0.50),
             tech_text, size=8, color=INK)
    cur_y += Inches(0.50)

    # Divider
    div = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
        body_x, int(cur_y), body_x + body_w, int(cur_y))
    div.line.color.rgb = RULE; div.line.width = Pt(0.5)
    cur_y += Inches(0.10)

    # Business / sales-closure context
    add_text(slide, body_x, int(cur_y), body_w, Inches(0.18),
             "BUSINESS CONTEXT", size=8, bold=True, color=integ["color"])
    cur_y += Inches(0.18)
    remaining = (card_top + card_h) - cur_y - Inches(0.1)
    add_text(slide, body_x, int(cur_y), body_w, int(remaining),
             integ["business"], size=9, color=INK)

# --- Footer ---
add_text(slide, Inches(0.5), SH - Inches(0.38), SW - Inches(2.0), Inches(0.3),
         "DealPad Integrations Architecture  •  Pricing & Scoping 2.0",
         size=9, color=MUTED)
add_text(slide, SW - Inches(1.6), SH - Inches(0.38), Inches(1.1), Inches(0.3),
         "1 / 1", size=9, color=MUTED, align=PP_ALIGN.RIGHT)

# ---------- write out ----------
out_dir = "exports"
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "DealPad_Integrations_Architecture.pptx")
prs.save(out_path)
print(out_path)
